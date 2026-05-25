"""
Drawing Primitives — مذكرتي Pro v17.1
Low-level, deterministic shape/text builders.

FIXED BUGS vs v17.0:
- gradient_fill() now inserts gradFill in CORRECT OOXML position (before <a:ln>)
- shadow() now inserts effectLst in CORRECT position (after <a:ln>)
- _sort_spPr() enforces strict OOXML child order on every shape
- set_solid_alpha() exposed (no duplication with slides.py)
"""
from __future__ import annotations

from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# Slide dimensions (cm) — 16:9
W, H = 33.867, 19.05

# OOXML spPr required child order
_SPPR_ORDER = [
    qn('a:xfrm'), qn('a:prstGeom'), qn('a:custGeom'),
    qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill'),
    qn('a:blipFill'), qn('a:pattFill'), qn('a:grpFill'),
    qn('a:ln'),
    qn('a:effectLst'), qn('a:effectDag'),
    qn('a:scene3d'), qn('a:sp3d'), qn('a:extLst'),
]
_SPPR_RANK = {tag: i for i, tag in enumerate(_SPPR_ORDER)}


def _sort_spPr(spPr) -> None:
    """Reorder <p:spPr> children to comply with OOXML schema."""
    children = list(spPr)
    children.sort(key=lambda el: _SPPR_RANK.get(el.tag, 99))
    for child in children:
        spPr.remove(child)
    for child in children:
        spPr.append(child)


def _get_spPr(shape):
    return shape._element.find(qn('p:spPr'))


# ── Unit helpers ─────────────────────────────────────────────────────
def cm(v: float) -> int:
    return int(Cm(v))

def pt(v: float) -> int:
    return int(Pt(v))


# ── Shape builders ───────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill: RGBColor, line=None, line_w=0.5):
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = pt(line_w)
    else:
        s.line.fill.background()
    return s


def rrect(slide, x, y, w, h, fill: RGBColor, radius_pct=8, line=None, line_w=0.5):
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(5, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = pt(line_w)
    else:
        s.line.fill.background()
    try:
        adj = s.adjustments
        if adj and len(adj) > 0:
            adj[0] = max(0, min(50, radius_pct)) * 1000
    except Exception:
        pass
    return s


def oval(slide, x, y, w, h, fill: RGBColor, alpha=100):
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(9, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def bg(slide, color: RGBColor):
    rect(slide, 0, 0, W, H, color)

def hline(slide, x, y, w, color: RGBColor, thickness=0.08):
    rect(slide, x, y, w, thickness, color)

def vline(slide, x, y, h2, color: RGBColor, thickness=0.08):
    rect(slide, x, y, thickness, h2, color)


# ── XML fill helpers ─────────────────────────────────────────────────
def set_solid_alpha(shape, alpha_pct: int) -> None:
    """Set transparency on a solidFill shape (0=transparent, 100=opaque)."""
    try:
        spPr = _get_spPr(shape)
        srgb = spPr.find('.//' + qn('a:srgbClr'))
        if srgb is not None:
            for e in srgb.findall(qn('a:alpha')):
                srgb.remove(e)
            alp = etree.SubElement(srgb, qn('a:alpha'))
            alp.set('val', str(int(alpha_pct * 1000)))
    except Exception:
        pass


def gradient_fill(shape, c1: str, c2: str, angle: float = 90) -> None:
    """
    Apply linear gradient via XML, with correct OOXML element ordering.
    BUG FIX: previously appended gradFill AFTER <a:ln> — now sorted correctly.
    """
    try:
        spPr = _get_spPr(shape)
        # Remove all existing fill variants
        for tag in [qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'),
                    qn('a:pattFill'), qn('a:blipFill'), qn('a:grpFill')]:
            for el in spPr.findall(tag):
                spPr.remove(el)

        # Build gradFill
        grad = etree.Element(qn('a:gradFill'))
        gsLst = etree.SubElement(grad, qn('a:gsLst'))

        gs0 = etree.SubElement(gsLst, qn('a:gs'))
        gs0.set('pos', '0')
        etree.SubElement(gs0, qn('a:srgbClr')).set('val', c1.lstrip('#'))

        gs1 = etree.SubElement(gsLst, qn('a:gs'))
        gs1.set('pos', '100000')
        etree.SubElement(gs1, qn('a:srgbClr')).set('val', c2.lstrip('#'))

        lin = etree.SubElement(grad, qn('a:lin'))
        lin.set('ang', str(int(angle * 60000)))
        lin.set('scaled', '0')

        spPr.append(grad)
        _sort_spPr(spPr)  # ← enforce correct order
    except Exception:
        pass


def gradient_rect(slide, x, y, w, h, c1: str, c2: str, angle=0):
    c1h = c1.lstrip('#')
    fill_color = RGBColor(int(c1h[0:2], 16), int(c1h[2:4], 16), int(c1h[4:6], 16))
    s = rect(slide, x, y, w, h, fill_color)
    if s:
        gradient_fill(s, c1, c2, angle)
    return s


def shadow(shape, blur=16, dist=5, angle=135, alpha=0.22, color="000000") -> None:
    """
    Add outer drop shadow via XML.
    BUG FIX: effectLst now inserted in correct OOXML position (after <a:ln>).
    """
    try:
        spPr = _get_spPr(shape)
        for old in spPr.findall(qn('a:effectLst')):
            spPr.remove(old)

        eLst = etree.Element(qn('a:effectLst'))
        shdw = etree.SubElement(eLst, qn('a:outerShdw'))
        shdw.set('blurRad', str(int(blur * 12700)))
        shdw.set('dist', str(int(dist * 12700)))
        shdw.set('dir', str(int(angle * 60000)))
        shdw.set('algn', 'tl')
        srgb = etree.SubElement(shdw, qn('a:srgbClr'))
        srgb.set('val', color.lstrip('#'))
        alp = etree.SubElement(srgb, qn('a:alpha'))
        alp.set('val', str(int(alpha * 100000)))

        spPr.append(eLst)
        _sort_spPr(spPr)  # ← enforce correct order
    except Exception:
        pass


# ── Text ─────────────────────────────────────────────────────────────
def txt(slide, text: str, x, y, w, h,
        font="Cairo", size=14, bold=False, italic=False,
        color: RGBColor | None = None,
        align=PP_ALIGN.RIGHT,
        margin=0.1, rtl=True, spacing=None,
        vcenter=True, line_spacing=1.15):
    """
    نص احترافي باستخدام Shape مع توسيط عمودي حقيقي.
    vcenter=True → MSO_ANCHOR.MIDDLE (يعمل في PowerPoint وLibreOffice)
    line_spacing → ارتفاع السطر النسبي
    """
    if not text or w <= 0 or h <= 0:
        return None
    try:
        from pptx.enum.text import MSO_ANCHOR
        sh = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
        sh.fill.background()
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left   = cm(margin)
        tf.margin_right  = cm(margin)
        tf.margin_top    = cm(0.04)
        tf.margin_bottom = cm(0.04)

        if vcenter:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            tf.vertical_anchor = MSO_ANCHOR.TOP

        p = tf.paragraphs[0]
        p.alignment = align

        # ارتفاع السطر
        try:
            from pptx.oxml.ns import qn as _qn
            pPr = p._p.get_or_add_pPr()
            # أزل lnSpc القديم إن وُجد
            for old in pPr.findall(_qn('a:lnSpc')):
                pPr.remove(old)
            lnSpc = etree.SubElement(pPr, _qn('a:lnSpc'))
            spcPct = etree.SubElement(lnSpc, _qn('a:spcPct'))
            spcPct.set('val', str(int(line_spacing * 100000)))
        except Exception:
            pass

        run = p.add_run()
        run.text = str(text)
        run.font.name   = font
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return sh
    except Exception:
        # fallback إلى textbox
        tb = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
        tb.word_wrap = True
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.name   = font
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return tb


def txt2(slide, label: str, value: str, x, y, w, h,
         font="Cairo", label_size=10, value_size=13,
         label_color: RGBColor | None = None,
         value_color: RGBColor | None = None,
         align=PP_ALIGN.RIGHT, margin=0.1):
    """
    نص بسطرين: تسمية + قيمة مع توسيط عمودي.
    مثالي لبطاقات المعلومات.
    """
    if w <= 0 or h <= 0: return None
    try:
        from pptx.enum.text import MSO_ANCHOR
        sh = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
        sh.fill.background()
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left   = cm(margin)
        tf.margin_right  = cm(margin)
        tf.margin_top    = cm(0.04)
        tf.margin_bottom = cm(0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p1 = tf.paragraphs[0]
        p1.alignment = align
        r1 = p1.add_run()
        r1.text = str(label)
        r1.font.name  = font
        r1.font.size  = Pt(label_size)
        r1.font.bold  = True
        if label_color: r1.font.color.rgb = label_color

        p2 = tf.add_paragraph()
        p2.alignment = align
        r2 = p2.add_run()
        r2.text = str(value)
        r2.font.name  = font
        r2.font.size  = Pt(value_size)
        r2.font.bold  = False
        if value_color: r2.font.color.rgb = value_color
        return sh
    except Exception:
        return None


def blank_slide(prs):
    """Add a completely blank slide (layout 6 = Blank)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Advanced Visual Primitives ────────────────────────────────────────

def glow(shape, color: str = "C6A03C", radius: int = 20, alpha: float = 0.4) -> None:
    """Add a glow effect (outerShdw with zero distance = glow)."""
    try:
        spPr = _get_spPr(shape)
        for old in spPr.findall(qn('a:effectLst')):
            spPr.remove(old)
        eLst = etree.Element(qn('a:effectLst'))
        g = etree.SubElement(eLst, qn('a:outerShdw'))
        g.set('blurRad', str(int(radius * 12700)))
        g.set('dist', '0')
        g.set('dir', '0')
        g.set('algn', 'ctr')
        srgb = etree.SubElement(g, qn('a:srgbClr'))
        srgb.set('val', color.lstrip('#'))
        alp = etree.SubElement(srgb, qn('a:alpha'))
        alp.set('val', str(int(alpha * 100000)))
        spPr.append(eLst)
        _sort_spPr(spPr)
    except Exception:
        pass


def multi_stop_gradient(shape, stops: list[tuple[int, str]], angle: float = 90) -> None:
    """
    Apply a multi-stop gradient.
    stops = [(pos_pct, '#RRGGBB'), ...]  e.g. [(0,'#07172F'),(50,'#1A3A6E'),(100,'#07172F')]
    """
    try:
        spPr = _get_spPr(shape)
        for tag in [qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'),
                    qn('a:pattFill'), qn('a:blipFill'), qn('a:grpFill')]:
            for el in spPr.findall(tag):
                spPr.remove(el)
        grad = etree.Element(qn('a:gradFill'))
        gsLst = etree.SubElement(grad, qn('a:gsLst'))
        for pos_pct, hex_color in stops:
            gs = etree.SubElement(gsLst, qn('a:gs'))
            gs.set('pos', str(int(pos_pct * 1000)))
            etree.SubElement(gs, qn('a:srgbClr')).set('val', hex_color.lstrip('#'))
        lin = etree.SubElement(grad, qn('a:lin'))
        lin.set('ang', str(int(angle * 60000)))
        lin.set('scaled', '0')
        spPr.append(grad)
        _sort_spPr(spPr)
    except Exception:
        pass


def gradient_oval(slide, x, y, w, h, c1: str, c2: str, angle=90, alpha=100):
    """Oval with gradient fill."""
    if w <= 0 or h <= 0:
        return None
    c1h = c1.lstrip('#')
    fill_color = RGBColor(int(c1h[0:2], 16), int(c1h[2:4], 16), int(c1h[4:6], 16))
    s = slide.shapes.add_shape(9, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    gradient_fill(s, c1, c2, angle)
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def triangle(slide, x, y, w, h, fill: RGBColor, pointing='up'):
    """Equilateral-ish triangle shape."""
    if w <= 0 or h <= 0:
        return None
    # Use right-triangle preset (shape 6) rotated for pointing direction
    shape_id = 6  # rtTriangle
    s = slide.shapes.add_shape(shape_id, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if pointing == 'down':
        s.rotation = 180
    elif pointing == 'left':
        s.rotation = 90
    elif pointing == 'right':
        s.rotation = 270
    return s


def diamond(slide, x, y, w, h, fill: RGBColor, alpha=100):
    """Diamond shape."""
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(4, cm(x), cm(y), cm(w), cm(h))  # shape 4 = diamond
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def hexagon(slide, x, y, w, h, fill: RGBColor, alpha=100):
    """Hexagon shape."""
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(10, cm(x), cm(y), cm(w), cm(h))  # shape 10 = hexagon
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def arc_progress(slide, x, y, size, fill: RGBColor, bg_color: RGBColor,
                 thickness=0.4) -> None:
    """Simulated progress ring using two arcs (outer ring + inner mask)."""
    # Outer ring background
    outer = oval(slide, x, y, size, size, bg_color, alpha=30)
    # Inner circle mask (creates ring effect)
    inner_offset = thickness
    inner_s = size - 2 * inner_offset
    oval(slide, x + inner_offset, y + inner_offset, inner_s, inner_s, bg_color, alpha=80)


def decorative_dots(slide, x, y, cols, rows, dot_size, gap, color: RGBColor, alpha=20):
    """Grid of decorative dots."""
    for r in range(rows):
        for c in range(cols):
            dx = x + c * (dot_size + gap)
            dy = y + r * (dot_size + gap)
            o = oval(slide, dx, dy, dot_size, dot_size, color)
            if o and alpha < 100:
                set_solid_alpha(o, alpha)


def wave_rect(slide, x, y, w, h, fill: RGBColor, wavy_top=True):
    """Rectangle with rounded top (simulates wave). Uses rrect with high radius."""
    if wavy_top:
        return rrect(slide, x, y, w, h, fill, radius_pct=12)
    return rect(slide, x, y, w, h, fill)


def badge(slide, x, y, w, h, fill_c1: str, fill_c2: str, label: str,
          font="Cairo", font_size=11, text_color: RGBColor = None, T=None):
    """Pill-shaped badge with gradient and centered label."""
    b = rrect(slide, x, y, w, h, RGBColor(0xC6, 0xA0, 0x3C), radius_pct=50)
    if b:
        gradient_fill(b, fill_c1, fill_c2, angle=0)
    if text_color is None and T is not None:
        text_color = T.text_dark_rgb
    txt(slide, label, x, y, w, h,
        font=font, size=font_size, bold=True,
        color=text_color, align=PP_ALIGN.CENTER, rtl=True)
    return b


def icon_circle(slide, x, y, size, bg_c1: str, bg_c2: str,
                icon: str, icon_size=20, T=None):
    """Circle with gradient bg + centered emoji/icon."""
    c = oval(slide, x, y, size, size,
             RGBColor(int(bg_c1.lstrip('#')[0:2], 16),
                      int(bg_c1.lstrip('#')[2:4], 16),
                      int(bg_c1.lstrip('#')[4:6], 16)))
    if c:
        gradient_fill(c, bg_c1, bg_c2, angle=135)
    txt(slide, icon, x, y, size, size,
        font="Calibri", size=icon_size, bold=False,
        color=T.text_dark_rgb if T else RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, rtl=False)
    return c


def number_badge(slide, x, y, size, num: int | str, T):
    """Circular number badge with accent gradient."""
    c = oval(slide, x, y, size, size, T.accent_rgb)
    if c:
        gradient_fill(c, T.accent_grad1, T.accent_grad2, 135)
        shadow(c, blur=10, dist=3, alpha=0.35)
    txt(slide, str(num), x, y, size, size,
        font="Calibri", size=max(8, int(size * 18)), bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False)
    return c


def divider(slide, x, y, w, T, style='gradient'):
    """Decorative divider line."""
    if style == 'gradient':
        d = rect(slide, x, y, w, 0.06, T.accent_rgb)
        if d:
            multi_stop_gradient(d, [(0, T.bg), (50, T.accent), (100, T.bg)], angle=0)
    elif style == 'double':
        rect(slide, x, y, w, 0.05, T.accent_rgb)
        rect(slide, x + w * 0.05, y + 0.12, w * 0.9, 0.03, T.muted_rgb)
    else:
        rect(slide, x, y, w, 0.06, T.accent_rgb)


def card_3d(slide, x, y, w, h, T, radius=10):
    """Card with shadow + subtle gradient for 3D feel."""
    # Shadow layer (slightly offset)
    shadow_s = rrect(slide, x + 0.15, y + 0.2, w, h, T.bg_rgb, radius_pct=radius)
    if shadow_s:
        set_solid_alpha(shadow_s, 40)

    # Main card
    c = rrect(slide, x, y, w, h, T.card_rgb, radius_pct=radius)
    if c:
        multi_stop_gradient(c, [(0, T.card), (100, T.bg2)], angle=135)
        shadow(c, blur=18, dist=5, alpha=0.4)
    return c


def slide_number(slide, num: int, total: int, T):
    """Slide number indicator bottom-right."""
    label = f"{num} / {total}"
    txt(slide, label, W - 3.5, H - 0.55, 3.2, 0.45,
        font="Calibri", size=9, bold=False,
        color=T.muted_rgb, align=PP_ALIGN.LEFT, rtl=False)


def watermark(slide, text: str, T):
    """Subtle watermark bottom-left."""
    txt(slide, text, 0.4, H - 0.55, 6.0, 0.45,
        font="Calibri", size=8, bold=False,
        color=T.muted_rgb, align=PP_ALIGN.RIGHT, rtl=False)


def section_tag(slide, label: str, x, y, T):
    """Small colored tag/label."""
    w, h = 3.5, 0.52
    b = rrect(slide, x, y, w, h, T.accent_rgb, radius_pct=50)
    if b:
        gradient_fill(b, T.accent_grad1, T.accent_grad2, 0)
    txt(slide, label, x, y, w, h,
        font="Cairo", size=10, bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=True)
