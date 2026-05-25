"""
Export Pipeline — مذكرتي Pro v17.1
Single entry point. Deterministic. No threads. No disk I/O.

FIXED vs v17.0:
- slide type set to 'custom' (fixes screen4x3 mismatch)
- per-slide tracing with fail-fast on any exception  
- full ZIP integrity validation after save
- Arabic filename sanitizer
"""
from __future__ import annotations

import io
import logging
import os
import shutil
import subprocess
import time
import zipfile
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Cm
from pptx.oxml.ns import qn

from core.models import PresentationRequest
from core.themes import get_theme
import engine.slides as _slides_canva
import engine.slides_premium as _slides_premium
import engine.slides_classic as _slides_classic

# خريطة المحركات
_ENGINES = {
    "canva":   _slides_canva,
    "premium": _slides_premium,
    "classic": _slides_classic,
}

def set_font(font_name: str):
    for mod in _ENGINES.values():
        mod.set_font(font_name)

def _get_fns(engine: str):
    mod = _ENGINES.get(engine, _slides_canva)
    return (
        mod.make_cover, mod.make_intro, mod.make_plan, mod.make_problem,
        mod.make_objectives, mod.make_importance, mod.make_methodology,
        mod.make_stats, mod.make_results, mod.make_conclusion,
        mod.make_recommendations, mod.make_future, mod.make_references,
        mod.make_final,
    )

log = logging.getLogger(__name__)
W_CM, H_CM = 33.867, 19.05
MIN_BYTES = 8_000


def _detect_font() -> str:
    candidates = ["Cairo", "Amiri", "Tahoma", "Arial Unicode MS", "Calibri"]
    if shutil.which("fc-list"):
        try:
            out = subprocess.run(
                ["fc-list", "--format=%{family}\n"],
                capture_output=True, text=True, timeout=5
            ).stdout.lower()
            for f in candidates:
                if f.lower() in out:
                    log.info(f"Font: {f}")
                    return f
        except Exception:
            pass
    font_dirs = ["/usr/share/fonts", "/usr/local/share/fonts",
                 os.path.expanduser("~/.fonts"), "/tmp/fonts"]
    for font in candidates[:3]:
        for d in font_dirs:
            if not os.path.isdir(d):
                continue
            for root, _, files in os.walk(d):
                for f in files:
                    if font.lower() in f.lower() and f.lower().endswith((".ttf", ".otf")):
                        log.info(f"Font (disk): {font}")
                        return font
    log.warning("No Arabic font — Calibri fallback")
    return "Calibri"


def _fix_slide_type(prs: Presentation) -> None:
    """Remove 'type' attr from sldSz so apps treat it as custom (not screen4x3)."""
    try:
        sldSz = prs.element.find(qn('p:sldSz'))
        if sldSz is not None and 'type' in sldSz.attrib:
            del sldSz.attrib['type']
    except Exception:
        pass


@dataclass
class ExportResult:
    success: bool
    data: bytes = b""
    slide_count: int = 0
    font_used: str = ""
    error: str = ""
    stages: list = field(default_factory=list)
    elapsed: float = 0.0


class PPTXExportPipeline:

    def __init__(self):
        self._font = _detect_font()
        set_font(self._font)
        log.info(f"Pipeline ready | font={self._font} | engines={list(_ENGINES.keys())}")

    def build(self, req: PresentationRequest) -> ExportResult:
        t0 = time.monotonic()
        stages = []
        try:
            # Stage 1: validate
            stages.append("validate")
            errors = req.validate()
            if errors:
                return ExportResult(success=False, error=" | ".join(errors), stages=stages)

            # Stage 2: init presentation
            stages.append("init_prs")
            prs = Presentation()
            prs.slide_width = Cm(W_CM)
            prs.slide_height = Cm(H_CM)
            _fix_slide_type(prs)

            # Stage 3: load theme
            stages.append("load_theme")
            theme = get_theme(req.theme)

            # Stage 4: build slides (per-slide tracing)
            stages.append("build_slides")
            n = self._build_slides(prs, req, theme, stages, req.engine)

            # Stage 5: serialize
            stages.append("serialize")
            data = self._serialize(prs)

            # Stage 6: validate output
            stages.append("validate_output")
            self._validate(data, n)

            elapsed = time.monotonic() - t0
            log.info(f"OK | slides={n} theme={req.theme} {len(data):,}B {elapsed:.2f}s")
            return ExportResult(success=True, data=data, slide_count=n,
                                font_used=self._font, stages=stages, elapsed=elapsed)

        except Exception as exc:
            stage = stages[-1] if stages else "unknown"
            log.error(f"FAIL [{stage}]: {exc}", exc_info=True)
            return ExportResult(success=False, error=f"[{stage}] {exc}",
                                stages=stages, elapsed=time.monotonic() - t0)

    def _build_slides(self, prs, req, theme, stages, engine="canva") -> int:
        (make_cover, make_intro, make_plan, make_problem,
         make_objectives, make_importance, make_methodology,
         make_stats, make_results, make_conclusion,
         make_recommendations, make_future, make_references,
         make_final) = _get_fns(engine)

        log.info(f"Engine: {engine}")
        cfg = req.slides
        count = 0

        # حساب العدد الحقيقي للشرائح قبل البناء
        def _will_run(condition): return bool(condition)
        total_slides = sum([
            _will_run(True),
            _will_run(cfg.intro and bool(req.intro_overview or req.intro_approach)),
            _will_run(cfg.plan and bool(req.chapters)),
            _will_run(cfg.problem and bool(req.main_problem or req.main_question or req.sub_questions)),
            _will_run(cfg.objectives and bool(req.objectives or req.hypotheses)),
            _will_run(cfg.importance and bool(req.importance or req.reasons)),
            _will_run(cfg.methodology and bool(req.methodology or req.sample_type or req.tool)),
            _will_run(cfg.kpi and bool(req.stats)),
            _will_run(cfg.results and bool(req.main_results)),
            _will_run(cfg.conclusion and bool(req.general_conclusion)),
            _will_run(cfg.recommendations and bool(req.recommendations)),
            _will_run(cfg.future and bool(req.future_work)),
            _will_run(cfg.references and bool(req.references)),
            _will_run(cfg.thankyou),
        ])
        # نمرر العدد الكلي للـ req ليستخدمه المحرك
        req._total_slides = total_slides

        def run(name, condition, fn):
            nonlocal count
            if not condition:
                return
            stages.append(f"slide:{name}")
            fn(prs, req, theme)
            count += 1
            log.debug(f"  slide:{name} OK")

        run("cover",           True,                                     make_cover)
        run("intro",           cfg.intro and bool(req.intro_overview or req.intro_approach), make_intro)
        run("plan",            cfg.plan and bool(req.chapters),          make_plan)
        run("problem",         cfg.problem and bool(req.main_problem or req.main_question or req.sub_questions), make_problem)
        run("objectives",      cfg.objectives and bool(req.objectives or req.hypotheses), make_objectives)
        run("importance",      cfg.importance and bool(req.importance or req.reasons),    make_importance)
        run("methodology",     cfg.methodology and bool(req.methodology or req.sample_type or req.tool), make_methodology)
        run("kpi",             cfg.kpi and bool(req.stats),              make_stats)
        run("results",         cfg.results and bool(req.main_results),   make_results)
        run("conclusion",      cfg.conclusion and bool(req.general_conclusion), make_conclusion)
        run("recommendations", cfg.recommendations and bool(req.recommendations), make_recommendations)
        run("future",          cfg.future and bool(req.future_work),     make_future)
        run("references",      cfg.references and bool(req.references),  make_references)
        run("thankyou",        cfg.thankyou,                             make_final)

        return count

    def _serialize(self, prs: Presentation) -> bytes:
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    def _validate(self, data: bytes, expected_slides: int) -> None:
        if len(data) < MIN_BYTES:
            raise ValueError(f"Output too small: {len(data)}B")
        if not data.startswith(b'PK'):
            raise ValueError("Not a valid ZIP/PPTX")
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                names = z.namelist()
                if 'ppt/presentation.xml' not in names:
                    raise ValueError("Missing ppt/presentation.xml")
                slides_in_zip = [n for n in names if n.startswith('ppt/slides/slide') and n.endswith('.xml')]
                if len(slides_in_zip) < expected_slides:
                    raise ValueError(f"Expected {expected_slides} slides, ZIP has {len(slides_in_zip)}")
                # Verify each slide is readable
                for s in slides_in_zip:
                    xml = z.read(s)
                    if len(xml) < 200:
                        raise ValueError(f"Slide too small: {s} ({len(xml)}B)")
        except zipfile.BadZipFile as e:
            raise ValueError(f"Corrupt ZIP: {e}")


_pipeline: PPTXExportPipeline | None = None


def get_pipeline() -> PPTXExportPipeline:
    global _pipeline
    if _pipeline is None:
        _pipeline = PPTXExportPipeline()
    return _pipeline
